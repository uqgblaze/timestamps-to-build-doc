"""
Build Document Generator — V2
Reads a *-timestamps.md and a .vtt file from a directory, then:
  V1 — produces a Markdown build document (each slide title = heading,
       matching transcript = paragraph body).
  V2 — optionally inserts the same content into each corresponding
       PowerPoint slide's speaker notes (title bold, then paragraph).
"""

import glob
import os
import re
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk


# ── time helpers ──────────────────────────────────────────────────────────────

def hms_to_ms(hms: str) -> int:
    m = re.match(r"^(\d{1,2}):(\d{2}):(\d{2})(?:\.(\d{1,3}))?$", hms.strip())
    if not m:
        raise ValueError(f"Invalid time: {hms!r}")
    frac = (m.group(4) or "0").ljust(3, "0")
    return (int(m.group(1)) * 3600 + int(m.group(2)) * 60 + int(m.group(3))) * 1000 + int(frac)


# ── file parsers ──────────────────────────────────────────────────────────────

def parse_timestamps(text: str) -> list[dict]:
    segments = []
    for line in text.splitlines():
        line = line.strip()
        if not line.startswith("|"):
            continue
        if "Start" in line or re.search(r":[-:]+", line):
            continue
        parts = [p.strip() for p in line.strip("|").split("|")]
        if len(parts) < 4:
            continue
        start, stop, title = parts[0], parts[2], parts[3].strip()
        if not re.match(r"^\d{1,2}:\d{2}:\d{2}(\.\d{1,3})?$", start):
            continue
        if not re.match(r"^\d{1,2}:\d{2}:\d{2}(\.\d{1,3})?$", stop):
            continue
        segments.append({"start": start, "stop": stop, "title": title})
    return segments


def parse_vtt(text: str) -> list[dict]:
    lines = text.replace("\ufeff", "").splitlines()
    timing_re = re.compile(
        r"^(\d{1,2}:\d{2}:\d{2}(?:\.\d{3})?)\s*-->\s*(\d{1,2}:\d{2}:\d{2}(?:\.\d{3})?)(?:\s+.*)?$"
    )
    cues, i = [], 0
    while i < len(lines) and not lines[i].strip().upper().startswith("WEBVTT"):
        i += 1
    i += 1
    while i < len(lines) and lines[i].strip():
        i += 1

    while i < len(lines):
        line = lines[i].strip()
        if not line:
            i += 1
            continue
        if not timing_re.match(line) and i + 1 < len(lines) and timing_re.match(lines[i + 1].strip()):
            i += 1
            line = lines[i].strip()
        m = timing_re.match(line)
        if not m:
            i += 1
            continue
        start_ms = hms_to_ms(m.group(1))
        end_ms   = hms_to_ms(m.group(2))
        i += 1
        text_lines = []
        while i < len(lines) and lines[i].strip():
            text_lines.append(lines[i].rstrip())
            i += 1
        raw = re.sub(r"<[^>]+>", "", " ".join(text_lines)).strip()
        if raw:
            cues.append({"start_ms": start_ms, "end_ms": end_ms, "text": raw})
    return cues


# ── core processing ───────────────────────────────────────────────────────────

def segment_slide_number(title: str) -> int | None:
    """Extract leading slide number from titles like '3. Exploration drill'."""
    m = re.match(r"^(\d+)\.", title.strip())
    return int(m.group(1)) if m else None


def build_segments_data(segments: list[dict], cues: list[dict]) -> list[dict]:
    """Return list of {slide_num, title, paragraph} for each segment."""
    data = []
    for idx, seg in enumerate(segments, start=1):
        start_ms = hms_to_ms(seg["start"])
        stop_ms  = hms_to_ms(seg["stop"])
        title    = seg["title"]
        slide_num = segment_slide_number(title) or idx

        matched = [
            c["text"] for c in cues
            if not (c["end_ms"] <= start_ms or c["start_ms"] >= stop_ms)
        ]
        paragraph = re.sub(r" {2,}", " ", " ".join(matched)) if matched else ""
        data.append({"slide_num": slide_num, "title": title, "paragraph": paragraph})
    return data


def build_markdown(data: list[dict], doc_title: str) -> str:
    lines = [f"# {doc_title}", ""]
    for seg in data:
        lines.append(f"## {seg['title']}")
        lines.append("")
        lines.append(seg["paragraph"] if seg["paragraph"] else "*(no transcript for this segment)*")
        lines.append("")
    return "\n".join(lines)


def insert_into_pptx(pptx_path: str, data: list[dict], out_path: str) -> list[str]:
    """Write title + paragraph into speaker notes for each matching slide.

    Returns a list of warning strings for slides that were out of range.
    """
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree

    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    warnings = []

    for seg in data:
        slide_num = seg["slide_num"]
        idx = slide_num - 1  # 0-based

        if idx < 0 or idx >= slide_count:
            warnings.append(
                f"Slide {slide_num} ({seg['title']!r}) — out of range "
                f"(deck has {slide_count} slides). Skipped."
            )
            continue

        slide = prs.slides[idx]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        txBody = tf._txBody

        # Remove existing paragraph elements, leaving bodyPr / lstStyle intact
        for p_elem in txBody.findall(qn("a:p")):
            txBody.remove(p_elem)

        def _add_para(text: str, bold: bool = False) -> None:
            p = etree.SubElement(txBody, qn("a:p"))
            if text:
                r = etree.SubElement(p, qn("a:r"))
                if bold:
                    rPr = etree.SubElement(r, qn("a:rPr"))
                    rPr.set("b", "1")
                t = etree.SubElement(r, qn("a:t"))
                t.text = text

        _add_para(seg["title"], bold=True)
        _add_para("")  # blank spacer
        _add_para(seg["paragraph"] or "(no transcript for this segment)")

    prs.save(out_path)
    return warnings


# ── file detection ────────────────────────────────────────────────────────────

def find_timestamps_files(directory: str) -> list[str]:
    patterns = ["*timestamps*.md", "*-timestamps*.md", "*_timestamps*.md"]
    found = []
    for pat in patterns:
        found.extend(glob.glob(os.path.join(directory, pat)))
    seen, unique = set(), []
    for f in found:
        key = os.path.normcase(f)
        if key not in seen:
            seen.add(key)
            unique.append(f)
    return unique or glob.glob(os.path.join(directory, "*.md"))


def find_vtt_files(directory: str) -> list[str]:
    return sorted(glob.glob(os.path.join(directory, "*.vtt")))


def find_pptx_files(directory: str) -> list[str]:
    return sorted(glob.glob(os.path.join(directory, "*.pptx")))


# ── GUI ───────────────────────────────────────────────────────────────────────

class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Build Document Generator — V2")
        self.resizable(True, True)
        self.minsize(680, 620)
        self._build_ui()
        self.directory = ""

    # ── layout ────────────────────────────────────────────────────────────────

    def _build_ui(self):
        pad = {"padx": 10, "pady": 4}

        # ── directory ──
        dir_frame = ttk.LabelFrame(self, text="1. Set Working Directory")
        dir_frame.pack(fill="x", **pad)
        self.dir_var = tk.StringVar()
        ttk.Entry(dir_frame, textvariable=self.dir_var, state="readonly").pack(
            side="left", fill="x", expand=True, padx=(6, 4), pady=6
        )
        ttk.Button(dir_frame, text="Browse…", command=self._browse).pack(
            side="left", padx=(0, 6), pady=6
        )

        # ── detected files ──
        files_frame = ttk.LabelFrame(self, text="2. Detected Files")
        files_frame.pack(fill="x", **pad)
        files_frame.columnconfigure(1, weight=1)

        ttk.Label(files_frame, text="Timestamps (.md):").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.ts_var = tk.StringVar()
        self.ts_combo = ttk.Combobox(files_frame, textvariable=self.ts_var, state="readonly", width=52)
        self.ts_combo.grid(row=0, column=1, sticky="ew", padx=6, pady=4)

        ttk.Label(files_frame, text="Transcript (.vtt):").grid(
            row=1, column=0, sticky="w", padx=6, pady=4
        )
        self.vtt_var = tk.StringVar()
        self.vtt_combo = ttk.Combobox(files_frame, textvariable=self.vtt_var, state="readonly", width=52)
        self.vtt_combo.grid(row=1, column=1, sticky="ew", padx=6, pady=4)

        ttk.Label(files_frame, text="Presentation (.pptx):").grid(
            row=2, column=0, sticky="w", padx=6, pady=4
        )
        self.pptx_var = tk.StringVar()
        self.pptx_combo = ttk.Combobox(files_frame, textvariable=self.pptx_var, state="readonly", width=52)
        self.pptx_combo.grid(row=2, column=1, sticky="ew", padx=6, pady=4)

        # ── V1 output ──
        v1_frame = ttk.LabelFrame(self, text="3(A). Markdown Output")
        v1_frame.pack(fill="x", **pad)
        v1_frame.columnconfigure(1, weight=1)

        ttk.Label(v1_frame, text="Document title:").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.title_var = tk.StringVar()
        ttk.Entry(v1_frame, textvariable=self.title_var).grid(
            row=0, column=1, columnspan=2, sticky="ew", padx=6, pady=4
        )

        ttk.Label(v1_frame, text="Output file (.md):").grid(
            row=1, column=0, sticky="w", padx=6, pady=4
        )
        self.md_out_var = tk.StringVar()
        ttk.Entry(v1_frame, textvariable=self.md_out_var).grid(
            row=1, column=1, sticky="ew", padx=6, pady=4
        )
        ttk.Button(v1_frame, text="Save as…", command=self._choose_md_output).grid(
            row=1, column=2, padx=(0, 6), pady=4
        )

        # ── V2 output ──
        v2_outer = ttk.LabelFrame(self, text="3(B). Insert into PowerPoint Speaker Notes (optional)")
        v2_outer.pack(fill="x", **pad)
        v2_outer.columnconfigure(0, weight=1)

        self.do_pptx_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            v2_outer,
            text="Also insert notes into the selected PowerPoint",
            variable=self.do_pptx_var,
            command=self._toggle_pptx_fields,
        ).grid(row=0, column=0, columnspan=3, sticky="w", padx=6, pady=(6, 2))

        self.v2_fields_frame = ttk.Frame(v2_outer)
        self.v2_fields_frame.grid(row=1, column=0, columnspan=3, sticky="ew")
        self.v2_fields_frame.columnconfigure(1, weight=1)

        ttk.Label(self.v2_fields_frame, text="Output file (.pptx):").grid(
            row=0, column=0, sticky="w", padx=6, pady=4
        )
        self.pptx_out_var = tk.StringVar()
        self.pptx_out_entry = ttk.Entry(self.v2_fields_frame, textvariable=self.pptx_out_var)
        self.pptx_out_entry.grid(row=0, column=1, sticky="ew", padx=6, pady=4)
        self.pptx_out_btn = ttk.Button(
            self.v2_fields_frame, text="Save as…", command=self._choose_pptx_output
        )
        self.pptx_out_btn.grid(row=0, column=2, padx=(0, 6), pady=4)

        self._toggle_pptx_fields()  # start disabled

        # ── generate ──
        btn_frame = ttk.Frame(self)
        btn_frame.pack(fill="x", padx=10, pady=6)
        ttk.Button(btn_frame, text="Generate", command=self._generate).pack(side="right")

        # ── log ──
        log_frame = ttk.LabelFrame(self, text="Log")
        log_frame.pack(fill="both", expand=True, **pad)
        self.log = scrolledtext.ScrolledText(
            log_frame, height=10, state="disabled", font=("Consolas", 9)
        )
        self.log.pack(fill="both", expand=True, padx=4, pady=4)

    # ── helpers ───────────────────────────────────────────────────────────────

    def _toggle_pptx_fields(self):
        state = "normal" if self.do_pptx_var.get() else "disabled"
        self.pptx_out_entry.configure(state=state)
        self.pptx_out_btn.configure(state=state)

    def _browse(self):
        d = filedialog.askdirectory(title="Select working directory")
        if not d:
            return
        self.directory = d
        self.dir_var.set(d)
        self._detect_files(d)

    def _detect_files(self, d: str):
        ts_files   = find_timestamps_files(d)
        vtt_files  = find_vtt_files(d)
        pptx_files = find_pptx_files(d)

        ts_names   = [os.path.basename(f) for f in ts_files]
        vtt_names  = [os.path.basename(f) for f in vtt_files]
        pptx_names = [os.path.basename(f) for f in pptx_files]

        self.ts_combo["values"]   = ts_names
        self.vtt_combo["values"]  = vtt_names
        self.pptx_combo["values"] = pptx_names

        if ts_names:
            self.ts_combo.current(0)
        if vtt_names:
            self.vtt_combo.current(0)
        if pptx_names:
            self.pptx_combo.current(0)

        stem = ""
        if vtt_names:
            stem = os.path.splitext(vtt_names[0])[0]
        elif ts_names:
            stem = re.sub(r"[-_]?timestamps?", "", os.path.splitext(ts_names[0])[0],
                          flags=re.IGNORECASE).strip("-_ ")

        self.title_var.set(stem)
        self.md_out_var.set(os.path.join(d, f"{stem}-build-doc.md") if stem else "")

        if pptx_names:
            base = os.path.splitext(pptx_names[0])[0]
            self.pptx_out_var.set(os.path.join(d, f"{base}-with-notes.pptx"))

        self._log(
            f"Found {len(ts_files)} timestamps, {len(vtt_files)} VTT, "
            f"{len(pptx_files)} PPTX file(s)."
        )

    def _choose_md_output(self):
        initial = self.md_out_var.get() or self.directory
        path = filedialog.asksaveasfilename(
            title="Save Markdown build document as",
            initialfile=os.path.basename(initial) if initial else "build-doc.md",
            initialdir=os.path.dirname(initial) if initial else self.directory,
            defaultextension=".md",
            filetypes=[("Markdown files", "*.md"), ("All files", "*.*")],
        )
        if path:
            self.md_out_var.set(path)

    def _choose_pptx_output(self):
        initial = self.pptx_out_var.get() or self.directory
        path = filedialog.asksaveasfilename(
            title="Save PowerPoint with notes as",
            initialfile=os.path.basename(initial) if initial else "presentation-with-notes.pptx",
            initialdir=os.path.dirname(initial) if initial else self.directory,
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
        )
        if path:
            self.pptx_out_var.set(path)

    # ── generate ──────────────────────────────────────────────────────────────

    def _generate(self):
        ts_name   = self.ts_var.get()
        vtt_name  = self.vtt_var.get()
        md_out    = self.md_out_var.get().strip()
        doc_title = self.title_var.get().strip() or "Build Document"
        do_pptx   = self.do_pptx_var.get()
        pptx_name = self.pptx_var.get()
        pptx_out  = self.pptx_out_var.get().strip()

        if not self.directory:
            messagebox.showerror("No directory", "Please select a working directory first.")
            return
        if not ts_name:
            messagebox.showerror("Missing file", "No timestamps file selected.")
            return
        if not vtt_name:
            messagebox.showerror("Missing file", "No VTT file selected.")
            return
        if not md_out:
            messagebox.showerror("No output path", "Please specify an output Markdown file path.")
            return
        if do_pptx:
            if not pptx_name:
                messagebox.showerror("Missing file", "No PowerPoint file selected.")
                return
            if not pptx_out:
                messagebox.showerror("No output path", "Please specify an output PPTX file path.")
                return

        ts_path   = os.path.join(self.directory, ts_name)
        vtt_path  = os.path.join(self.directory, vtt_name)
        pptx_path = os.path.join(self.directory, pptx_name) if pptx_name else ""

        self._log("─" * 60)
        self._log(f"Timestamps  : {ts_name}")
        self._log(f"Transcript  : {vtt_name}")
        self._log(f"MD output   : {md_out}")
        if do_pptx:
            self._log(f"PPTX source : {pptx_name}")
            self._log(f"PPTX output : {pptx_out}")

        # ── read sources ──
        try:
            with open(ts_path, encoding="utf-8") as f:
                ts_text = f.read()
            with open(vtt_path, encoding="utf-8-sig") as f:
                vtt_text = f.read()
        except OSError as e:
            messagebox.showerror("File error", str(e))
            return

        segments = parse_timestamps(ts_text)
        if not segments:
            messagebox.showerror("Parse error", "No valid segments found in the timestamps file.")
            return
        self._log(f"Parsed {len(segments)} segments.")

        cues = parse_vtt(vtt_text)
        self._log(f"Parsed {len(cues)} VTT cues.")

        data = build_segments_data(segments, cues)

        # ── V1: write markdown ──
        md_text = build_markdown(data, doc_title)
        try:
            with open(md_out, "w", encoding="utf-8") as f:
                f.write(md_text)
            self._log(f"Markdown written → {md_out}")
        except OSError as e:
            messagebox.showerror("Write error", str(e))
            return

        # ── V2: insert into PPTX ──
        if do_pptx:
            try:
                warnings = insert_into_pptx(pptx_path, data, pptx_out)
            except Exception as e:
                messagebox.showerror("PPTX error", str(e))
                return

            if warnings:
                self._log("Warnings:")
                for w in warnings:
                    self._log(f"  ⚠  {w}")
            written = len(data) - len(warnings)
            self._log(f"PPTX notes inserted for {written}/{len(data)} slides → {pptx_out}")

        # ── done ──
        summary = f"Markdown: {md_out}"
        if do_pptx:
            summary += f"\nPowerPoint: {pptx_out}"
        messagebox.showinfo("Done", f"Generated successfully.\n\n{summary}")

    def _log(self, msg: str):
        self.log.configure(state="normal")
        self.log.insert("end", msg + "\n")
        self.log.see("end")
        self.log.configure(state="disabled")


# ── entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    app = App()
    app.mainloop()
